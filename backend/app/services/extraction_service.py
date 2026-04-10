import io
import logging
from pathlib import Path

import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation

logger = logging.getLogger(__name__)


def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from a PDF file. Falls back to OCR if digital text is insufficient."""
    text_parts = []
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text() or ""
                text_parts.append(page_text)
    except Exception as e:
        logger.error(f"PDF extraction error: {e}")
        return ""

    full_text = "\n".join(text_parts).strip()

    # If very little text extracted, try OCR
    if len(full_text) < 100:
        return _ocr_pdf(file_path)

    return full_text


def _ocr_pdf(file_path: str) -> str:
    """OCR a scanned PDF using Tesseract."""
    try:
        from pdf2image import convert_from_path
        import pytesseract

        images = convert_from_path(file_path, dpi=300)
        text_parts = []
        for img in images:
            text = pytesseract.image_to_string(img, lang="eng")
            text_parts.append(text)
        return "\n".join(text_parts).strip()
    except Exception as e:
        logger.error(f"OCR error: {e}")
        return ""


def extract_text_from_docx(file_path: str) -> str:
    """Extract text from a Word document."""
    try:
        doc = DocxDocument(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n".join(paragraphs)
    except Exception as e:
        logger.error(f"DOCX extraction error: {e}")
        return ""


def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from a PowerPoint file."""
    try:
        prs = Presentation(file_path)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            text_parts.append(text)
        return "\n".join(text_parts)
    except Exception as e:
        logger.error(f"PPTX extraction error: {e}")
        return ""


def extract_text_from_txt(file_path: str) -> str:
    """Extract text from a plain text file."""
    try:
        return Path(file_path).read_text(encoding="utf-8")
    except Exception as e:
        logger.error(f"TXT extraction error: {e}")
        return ""


def extract_text(file_path: str, mime_type: str = "") -> str:
    """Extract text from a file based on its type."""
    path = file_path.lower()
    if path.endswith(".pdf") or "pdf" in mime_type:
        return extract_text_from_pdf(file_path)
    elif path.endswith(".docx") or "wordprocessingml" in mime_type:
        return extract_text_from_docx(file_path)
    elif path.endswith(".pptx") or "presentationml" in mime_type:
        return extract_text_from_pptx(file_path)
    elif path.endswith(".txt") or "text/plain" in mime_type:
        return extract_text_from_txt(file_path)
    else:
        logger.warning(f"Unsupported file type: {path}")
        return ""
